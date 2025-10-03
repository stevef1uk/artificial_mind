#!/usr/bin/env python3
"""
Generate Artificial Mind Solution Architecture PowerPoint deck.
Creates docs/Artificial_Mind_Solution_Architecture.pptx with sections:
  - Context & High-Level View
  - Component Deep Dives (next level down)
  - Technical Architecture
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, text in enumerate(bullets):
        p = body.add_paragraph() if i > 0 else body.paragraphs[0]
        p.text = text
        p.level = 0


def add_section_header(prs: Presentation, title: str) -> None:
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    # Slight styling
    run = title_shape.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)


def build_deck() -> Presentation:
    prs = Presentation()

    # 1) Title
    add_title_slide(
        prs,
        title="Artificial Mind Solution Architecture",
        subtitle="Context, High-Level View, Components, and Technical Architecture",
    )

    # 2) Context
    add_section_header(prs, "Context & Objectives")
    add_bullets_slide(
        prs,
        "Project Context",
        [
            "Artificial mind combining ethical decision-making, planning, and self-improvement",
            "Event-driven, service-oriented architecture for scalable cognition",
            "Safety-first execution with Principles gating and Docker sandboxing",
        ],
    )
    add_bullets_slide(
        prs,
        "Primary Objectives",
        [
            "Understand natural language tasks and generate executable solutions",
            "Learn and cache capabilities for faster, more reliable execution",
            "Maintain ethical boundaries and transparency",
        ],
    )

    # 3) High-Level Architecture
    add_section_header(prs, "High-Level Architecture")
    add_bullets_slide(
        prs,
        "Core Layers",
        [
            "FSM Engine: Orchestrates cognition (perceive → learn → plan → evaluate → execute)",
            "HDN: Intelligent code generation, validation, caching, execution",
            "Self-Model & Goal Manager: Goals, beliefs, episodes, prioritization",
            "Principles Server: Ethical/safety rules and pre-exec gating",
            "Event Bus (NATS): Canonical events for coordination and observability",
        ],
    )
    add_bullets_slide(
        prs,
        "Memory & Knowledge",
        [
            "Working Memory (Redis): Ephemeral state, capabilities, artifacts",
            "Episodic Memory (Qdrant): Vector search for execution history",
            "Semantic Knowledge (Neo4j): Domain concepts, relations, constraints",
        ],
    )

    # 4) Component Deep Dives
    add_section_header(prs, "Component Deep Dives (Next Level)")

    add_bullets_slide(
        prs,
        "FSM Engine",
        [
            "State-driven pipeline with reasoning, growth, and safety integration",
            "Delegates capability execution to HDN; consults goals/principles",
            "Curiosity goals from news and knowledge gaps",
        ],
    )
    add_bullets_slide(
        prs,
        "HDN Intelligent Execution",
        [
            "LLM-powered multi-language code generation",
            "Docker sandbox validation with retries and auto-fix",
            "Capability caching + dynamic action registration (Redis)",
        ],
    )
    add_bullets_slide(
        prs,
        "Principles Server",
        [
            "JSON rules, dynamic reload, context-aware checks",
            "Pre-exec gating for tools/actions, auditable denials",
        ],
    )
    add_bullets_slide(
        prs,
        "Self-Model & Goal Manager",
        [
            "Tracks goals (active/history/priorities), beliefs, episodes",
            "Publishes agi.goal.* on NATS; deduplication and scoring",
            "Influences planner/decider via active priorities",
        ],
    )
    add_bullets_slide(
        prs,
        "Planner / Evaluator",
        [
            "Generates multi-step workflows, ranks, executes, records episodes",
            "Consults domain knowledge and episodic memory for scoring",
        ],
    )
    add_bullets_slide(
        prs,
        "Monitor UI",
        [
            "Health, metrics, workflows, artifacts, capabilities",
            "Surfaces cold-start vs cached execution and safety status",
        ],
    )

    # 5) Technical Architecture
    add_section_header(prs, "Technical Architecture")
    add_bullets_slide(
        prs,
        "Technology Stack",
        [
            "Go services (FSM, HDN, Principles, Monitor)",
            "Redis (working memory, registry, artifacts)",
            "Qdrant (episodic memory), Neo4j (semantic knowledge)",
            "Docker for sandboxed execution; optional k3s manifests",
        ],
    )
    add_bullets_slide(
        prs,
        "Security & Safety",
        [
            "Principles pre-exec checks + LLM safety categorization",
            "Docker isolation with timeouts and resource limits",
            "Audit trails and metrics for tool usage and workflows",
        ],
    )
    add_bullets_slide(
        prs,
        "Scalability & Ops",
        [
            "Event-driven via NATS; stateless APIs for horizontal scaling",
            "Makefile for builds/tests; docker-compose for memory infra",
            "Monitor UI for status, metrics, and artifact access",
        ],
    )

    # 6) Summary / Close
    add_section_header(prs, "Summary")
    add_bullets_slide(
        prs,
        "Key Takeaways",
        [
            "Safety-conscious, self-improving architecture",
            "Modular services with layered memory and knowledge",
            "Operationally viable with tests, metrics, and UI",
        ],
    )

    return prs


def main() -> None:
    prs = build_deck()
    output = "docs/Artificial_Mind_Solution_Architecture.pptx"
    prs.save(output)
    print(f"✅ Generated {output}")


if __name__ == "__main__":
    main()


